"""
Export utilities for generating reports in various formats.
"""

import json
from typing import Dict, List, Any, Optional
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from docx import Document
from docx.shared import Inches, Pt, RGBColor


class ReportExporter:
    """Export reports in various formats."""
    
    def __init__(self, output_dir: str = "./outputs"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
    
    def export_json(self, data: Dict[str, Any], filename: str) -> str:
        """Export data to JSON file."""
        filepath = self.output_dir / f"{filename}.json"
        with open(filepath, 'w') as f:
            json.dump(data, f, indent=2, default=str)
        return str(filepath)
    
    def export_pdf(self, content: Dict[str, Any], filename: str) -> str:
        """Export report to PDF."""
        filepath = self.output_dir / f"{filename}.pdf"
        doc = SimpleDocTemplate(str(filepath), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30
        )
        story.append(Paragraph(content.get('title', 'Report'), title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Executive Summary
        if 'executive_summary' in content:
            story.append(Paragraph("Executive Summary", styles['Heading2']))
            story.append(Spacer(1, 0.1*inch))
            story.append(Paragraph(content['executive_summary'], styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        # Main Content
        if 'sections' in content:
            for section in content['sections']:
                story.append(Paragraph(section.get('title', ''), styles['Heading2']))
                story.append(Spacer(1, 0.1*inch))
                story.append(Paragraph(section.get('content', ''), styles['Normal']))
                story.append(Spacer(1, 0.2*inch))
        
        # Recommendations
        if 'recommendations' in content:
            story.append(Paragraph("Recommendations", styles['Heading2']))
            story.append(Spacer(1, 0.1*inch))
            for i, rec in enumerate(content['recommendations'], 1):
                story.append(Paragraph(f"{i}. {rec}", styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
        
        # Footer
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(
            f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            styles['Normal']
        ))
        
        doc.build(story)
        return str(filepath)
    
    def export_pptx(self, content: Dict[str, Any], filename: str) -> str:
        """Export presentation to PowerPoint."""
        filepath = self.output_dir / f"{filename}.pptx"
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = content.get('title', 'Report')
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Executive Summary slide
        if 'executive_summary' in content:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Executive Summary"
            tf = body_shape.text_frame
            tf.text = content['executive_summary']
        
        # Content slides
        if 'sections' in content:
            for section in content['sections']:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text = section.get('title', '')
                tf = body_shape.text_frame
                tf.text = section.get('content', '')
        
        # Recommendations slide
        if 'recommendations' in content:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Recommendations"
            tf = body_shape.text_frame
            for rec in content['recommendations']:
                p = tf.add_paragraph()
                p.text = rec
                p.level = 0
        
        prs.save(str(filepath))
        return str(filepath)
    
    def export_docx(self, content: Dict[str, Any], filename: str) -> str:
        """Export report to Word document."""
        filepath = self.output_dir / f"{filename}.docx"
        doc = Document()
        
        # Title
        title = doc.add_heading(content.get('title', 'Report'), 0)
        
        # Executive Summary
        if 'executive_summary' in content:
            doc.add_heading('Executive Summary', level=1)
            doc.add_paragraph(content['executive_summary'])
        
        # Main Content
        if 'sections' in content:
            for section in content['sections']:
                doc.add_heading(section.get('title', ''), level=1)
                doc.add_paragraph(section.get('content', ''))
        
        # Recommendations
        if 'recommendations' in content:
            doc.add_heading('Recommendations', level=1)
            for rec in content['recommendations']:
                doc.add_paragraph(rec, style='List Bullet')
        
        # Footer
        doc.add_paragraph(f"\nGenerated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        
        doc.save(str(filepath))
        return str(filepath)
    
    def export(
        self,
        content: Dict[str, Any],
        filename: str,
        format: str = "json"
    ) -> str:
        """
        Export content in specified format.
        
        Args:
            content: Content to export
            filename: Base filename (without extension)
            format: Export format (json, pdf, pptx, docx)
            
        Returns:
            Path to exported file
        """
        format = format.lower()
        if format == "json":
            return self.export_json(content, filename)
        elif format == "pdf":
            return self.export_pdf(content, filename)
        elif format == "pptx":
            return self.export_pptx(content, filename)
        elif format == "docx":
            return self.export_docx(content, filename)
        else:
            raise ValueError(f"Unsupported format: {format}")

